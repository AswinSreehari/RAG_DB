import sys
import os
import re
from pptx import Presentation
from PIL import Image, ImageFilter, ImageEnhance
import pytesseract
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()
tesseract_path = os.getenv("TESSERACT_PATH", r"C:\Program Files\Tesseract-OCR\tesseract.exe")
pytesseract.pytesseract.tesseract_cmd = tesseract_path

def preprocess_image(image_path):
    """Enhance image for better OCR accuracy."""
    img = Image.open(image_path).convert('L')  # Convert to grayscale
    img = img.filter(ImageFilter.SHARPEN)
    img = ImageEnhance.Contrast(img).enhance(2.0)
    return img

def clean_text(text):
    """Remove non-ASCII characters and extra whitespace."""
    text = re.sub(r'[^\x00-\x7F]+', ' ', text)  # Remove weird unicode
    text = re.sub(r'\s+', ' ', text)  # Collapse multiple spaces
    return text.strip()

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    full_text = []
    image_counter = 0
    img_dir = os.path.join(os.path.dirname(pptx_path), 'slide_images')
    os.makedirs(img_dir, exist_ok=True)

    for i, slide in enumerate(prs.slides):
        slide_text = [f"Slide {i+1}:"]

        for shape in slide.shapes:
            # Extract plain text
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(f"[Text]: {shape.text.strip()}")

            # Extract chart data
            if shape.has_chart:
                chart = shape.chart
                chart_block = ["[Chart]"]

                if chart.has_title and chart.chart_title.has_text_frame:
                    title_frame = chart.chart_title.text_frame
                    if title_frame and title_frame.text:
                        chart_block.append(f"Title: {title_frame.text.strip()}")

                for series in chart.series:
                    series_name = series.name or "Unnamed Series"
                    chart_block.append(f"Series: {series_name}")
                    try:
                        categories = [str(pt) for pt in series.category_axis_data]
                        values = [str(v) for v in series.values]
                        chart_block.append(f"  Categories: {', '.join(categories)}")
                        chart_block.append(f"  Values: {', '.join(values)}")
                    except Exception as e:
                        chart_block.append(f"  [Error extracting chart data: {e}]")

                slide_text.append("\n".join(chart_block))

            # OCR from images
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    img_ext = image.ext
                    img_bytes = image.blob
                    img_path = os.path.join(img_dir, f"slide{i+1}_img{image_counter}.{img_ext}")
                    
                    with open(img_path, "wb") as f:
                        f.write(img_bytes)

                    processed_img = preprocess_image(img_path)
                    ocr_text = pytesseract.image_to_string(processed_img, lang='eng', config='--psm 6')
                    if ocr_text.strip():
                        slide_text.append("[Image OCR]: " + clean_text(ocr_text))

                except Exception as e:
                    slide_text.append(f"[Image OCR error]: {e}")
                
                image_counter += 1

        full_text.append("\n".join(slide_text))

    return "\n\n".join(full_text)

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: pptx_text_extractor.py <path-to-pptx>", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    if not os.path.isfile(pptx_path):
        print(f"File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    result = extract_text_from_pptx(pptx_path)
    print(result)







# import sys
# import os
# from pptx import Presentation
# from PIL import Image, ImageFilter, ImageOps
# import pytesseract
# from dotenv import load_dotenv

# # Load environment variables from .env file
# load_dotenv()

# # Use TESSERACT_PATH from environment or fallback to default Windows install location
# tesseract_path = os.getenv("TESSERACT_PATH", r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe")
# pytesseract.pytesseract.tesseract_cmd = tesseract_path


# def preprocess_image(img_path):
#     img = Image.open(img_path).convert('L')  # Convert to grayscale
#     img = ImageOps.invert(img)               # Invert image (optional for light backgrounds)
#     img = img.point(lambda x: 0 if x < 150 else 255, '1')  # Apply binary thresholding
#     img = img.resize((img.width * 2, img.height * 2))  # Resize image to improve OCR accuracy
#     img = img.filter(ImageFilter.MedianFilter())  # Apply median filter to reduce noise
#     return img


# def extract_text_from_pptx(pptx_path):
#     prs = Presentation(pptx_path)
#     full_text = []
#     image_counter = 0
#     img_dir = os.path.join(os.path.dirname(pptx_path), 'slide_images')
#     os.makedirs(img_dir, exist_ok=True)

#     for i, slide in enumerate(prs.slides):
#         slide_text = [f"Slide {i+1}:"]

#         for shape in slide.shapes:
#             # Extract plain text
#             if hasattr(shape, "text") and shape.text.strip():
#                 slide_text.append(shape.text.strip())

#             # Extract chart data
#             if shape.has_chart:
#                 chart = shape.chart
#                 chart_block = ["[Chart]"]

#                 if chart.has_title and chart.chart_title.has_text_frame:
#                     title_frame = chart.chart_title.text_frame
#                     if title_frame and title_frame.text:
#                         chart_block.append(f"Title: {title_frame.text.strip()}")

#                 for series in chart.series:
#                     series_name = series.name or "Unnamed Series"
#                     chart_block.append(f"Series: {series_name}")
#                     try:
#                         categories = [pt for pt in series.category_axis_data]
#                         values = [pt for pt in series.values]
#                         chart_block.append(f"  Categories: {', '.join(str(c) for c in categories)}")
#                         chart_block.append(f"  Values: {', '.join(str(v) for v in values)}")
#                     except Exception as e:
#                         chart_block.append(f"  [Error extracting chart data: {e}]")

#                 slide_text.append("\n".join(chart_block))

#             # OCR from image shapes
#             if shape.shape_type == 13:  # Picture
#                 image = shape.image
#                 img_ext = image.ext
#                 img_bytes = image.blob
#                 img_path = os.path.join(img_dir, f"slide{i+1}_img{image_counter}.{img_ext}")
#                 with open(img_path, "wb") as f:
#                     f.write(img_bytes)
#                 try:
#                     clean_img = preprocess_image(img_path)
#                     ocr_text = pytesseract.image_to_string(clean_img, config='--psm 6')
#                     if ocr_text.strip():
#                         slide_text.append("[Image OCR]: " + ocr_text.strip())
#                 except Exception as e:
#                     slide_text.append(f"[Image OCR error]: {e}")
#                 image_counter += 1

#         full_text.append("\n".join(slide_text))

#     return "\n\n".join(full_text)


# if __name__ == "__main__":
#     if len(sys.argv) != 2:
#         print("Usage: pptx_text_extractor.py <path-to-pptx>", file=sys.stderr)
#         sys.exit(1)

#     pptx_path = sys.argv[1]
#     if not os.path.isfile(pptx_path):
#         print(f"File not found: {pptx_path}", file=sys.stderr)
#         sys.exit(1)

#     result = extract_text_from_pptx(pptx_path)
#     print(result)